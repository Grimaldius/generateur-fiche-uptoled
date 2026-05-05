#!/usr/bin/env python3
"""
generateur_fiche_uptoled.py
───────────────────────────
Convertit n'importe quel fichier fournisseur (JPG, PNG, PDF, PPTX)
en fiche technique Uptoled au format A3 paysage.

Usage:
    python3 generateur_fiche_uptoled.py <fichier_source> [output.pptx]

Prérequis:
    pip install anthropic PyMuPDF Pillow numpy python-pptx
    npm install -g pptxgenjs
    export ANTHROPIC_API_KEY=sk-ant-...
"""

import sys, os, json, base64, subprocess, shutil, tempfile
from pathlib import Path
import numpy as np

# ── Dépendances ──────────────────────────────────────────────────
try:
    from PIL import Image
    import fitz           # PyMuPDF
    import anthropic
except ImportError as e:
    print(f"[ERREUR] Dépendance manquante : {e}")
    print("Installe : pip install anthropic PyMuPDF Pillow numpy")
    sys.exit(1)

# ── Chemins fixes ────────────────────────────────────────────────
SCRIPT_DIR   = Path(__file__).parent
GENERATE_JS  = SCRIPT_DIR / "generate_pptx.js"
LOGO_PATH    = SCRIPT_DIR / "LogoUptoled_transparent.png"


# ════════════════════════════════════════════════════════════════
#  1. CONVERSION EN IMAGES
# ════════════════════════════════════════════════════════════════

def file_to_images(input_path: Path, tmp_dir: Path) -> list[Path]:
    """Convertit le fichier source en liste d'images PNG."""
    ext = input_path.suffix.lower()
    images = []

    if ext in (".jpg", ".jpeg", ".png"):
        # Copie directe — le fichier EST l'image
        dst = tmp_dir / f"page_1{ext}"
        shutil.copy(input_path, dst)
        images.append(dst)

    elif ext == ".pdf":
        doc = fitz.open(str(input_path))
        for i, page in enumerate(doc):
            mat = fitz.Matrix(2.0, 2.0)           # résolution ×2 pour lisibilité API
            pix = page.get_pixmap(matrix=mat)
            out = tmp_dir / f"page_{i+1}.png"
            pix.save(str(out))
            images.append(out)
        doc.close()

    elif ext in (".pptx", ".ppt"):
        # LibreOffice → PDF → images
        result = subprocess.run(
            ["python3", "/mnt/skills/public/pptx/scripts/office/soffice.py",
             "--headless", "--convert-to", "pdf", str(input_path)],
            capture_output=True, cwd=str(tmp_dir)
        )
        pdf_path = tmp_dir / (input_path.stem + ".pdf")
        if not pdf_path.exists():
            # LibreOffice écrit dans le répertoire courant
            fallback = Path(input_path.stem + ".pdf")
            if fallback.exists():
                shutil.move(str(fallback), str(pdf_path))
        if pdf_path.exists():
            images = file_to_images(pdf_path, tmp_dir)
        else:
            raise FileNotFoundError("Échec conversion PPTX → PDF")

    else:
        raise ValueError(f"Format non supporté : {ext}")

    return images


# ════════════════════════════════════════════════════════════════
#  2. EXTRACTION DONNÉES PRODUIT VIA CLAUDE API
# ════════════════════════════════════════════════════════════════

PROMPT_EXTRACTION = """
Tu analyses une fiche technique produit (éclairage LED professionnel).
Extrais TOUTES les informations disponibles et retourne UNIQUEMENT un JSON valide,
sans texte avant ni après, sans balises markdown.

Structure JSON attendue :
{
  "nom_produit": "NOM COMPLET DU PRODUIT EN MAJUSCULES",
  "reference": "REF-PRODUIT",
  "serie": "Éclairage industriel",
  "gamme": "Nom de gamme si présent",
  "accroche": "Une phrase courte bénéfice principal",
  "description": "2-3 phrases description complète",
  "specifications": [
    ["Matériau", "valeur"],
    ["Diffuseur", "valeur"],
    ["Raccordement", "valeur"],
    ["Presse-étoupe", "valeur"],
    ["Durée de vie", "valeur"],
    ["Temp. de couleur", "valeur"],
    ["IRC", "valeur"],
    ["Gradation", "valeur"],
    ["Montages", "valeur"],
    ["Origine", "valeur"],
    ["Garantie", "valeur"]
  ],
  "certifications": {
    "CE": true/false,
    "IP20": false, "IP40": false, "IP65": false, "IP66": false, "IP67": false, "IP69K": false,
    "IK08": false, "IK10": false,
    "IRC70": false, "IRC80": false, "IRC90": false,
    "DALI": false, "DSI": false, "V1_10": false, "V0_10": false, "TRIAC": false,
    "ATEX": false, "NF": false, "FAB_EU": false
  },
  "dimensions": [
    {"ref": "REF-XXX", "watts": "35", "flux": "4 470", "A": "703", "B": "310", "C": "142", "poids": "7,8"}
  ],
  "a_photo_produit": true/false,
  "a_diagramme_polaire": true/false,
  "a_schema_cote": true/false,
  "annee": "2026"
}

Règles :
- Si une spec n'est pas trouvée, mets null (pas "—")
- Pour les certifications, déduis depuis les marquages présents (ex: IP65 → IP65: true)
- Pour les dimensions : extrais TOUTES les lignes du tableau (W, Lm, A, B, C, poids)
- "a_photo_produit" = true si une photo du luminaire est visible
- "a_diagramme_polaire" = true si un diagramme polaire/courbe photométrique est visible
- "a_schema_cote" = true si un schéma coté avec cotes A/B/C est visible
"""

def extract_product_data(images: list[Path]) -> dict:
    """Envoie les images à Claude API et retourne les données structurées."""
    client = anthropic.Anthropic()

    content = []
    # On envoie au maximum 3 pages (évite token overflow)
    for img_path in images[:3]:
        with open(img_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        # Détermine le media_type
        mt = "image/png" if str(img_path).endswith(".png") else "image/jpeg"
        content.append({"type": "image",
                         "source": {"type": "base64", "media_type": mt, "data": b64}})

    content.append({"type": "text", "text": PROMPT_EXTRACTION})

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=2000,
        messages=[{"role": "user", "content": content}]
    )

    raw = response.content[0].text.strip()
    # Nettoyage des éventuelles balises markdown
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1].rsplit("```", 1)[0]

    return json.loads(raw)


# ════════════════════════════════════════════════════════════════
#  3. EXTRACTION IMAGE PRINCIPALE
# ════════════════════════════════════════════════════════════════

def extract_main_photo(input_path: Path, tmp_dir: Path) -> Path | None:
    """Extrait la photo principale du produit depuis le fichier source."""
    ext = input_path.suffix.lower()

    if ext in (".jpg", ".jpeg", ".png"):
        # Le fichier est lui-même la photo
        dst = tmp_dir / f"photo_produit{ext}"
        shutil.copy(input_path, dst)
        return dst

    elif ext == ".pdf":
        return _extract_largest_image_pdf(input_path, tmp_dir, prefix="photo")

    elif ext in (".pptx", ".ppt"):
        return _extract_largest_image_pptx(input_path, tmp_dir, prefix="photo")

    return None


def _extract_largest_image_pdf(pdf_path: Path, tmp_dir: Path, prefix: str) -> Path | None:
    """Extrait l'image la plus grande d'un PDF (heuristique = photo produit)."""
    doc   = fitz.open(str(pdf_path))
    best  = None
    best_area = 0

    for page in doc:
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            base_img = doc.extract_image(xref)
            w, h = base_img["width"], base_img["height"]
            area  = w * h
            # Filtre : ignore les images trop petites (icônes, logos <100×100)
            if area < 10000:
                continue
            if area > best_area:
                best_area = area
                best      = (xref, base_img["ext"])

    if best:
        xref, img_ext = best
        base_img = doc.extract_image(xref)
        out_path = tmp_dir / f"{prefix}_produit.{img_ext}"
        with open(out_path, "wb") as f:
            f.write(base_img["image"])
        doc.close()
        return out_path

    doc.close()
    return None


def _extract_largest_image_pptx(pptx_path: Path, tmp_dir: Path, prefix: str) -> Path | None:
    """Extrait l'image la plus grande d'un PPTX."""
    from pptx import Presentation as PptxPrs
    prs   = PptxPrs(str(pptx_path))
    best  = None
    best_area = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                img_bytes = shape.image.blob
                try:
                    img = Image.open(__import__("io").BytesIO(img_bytes))
                    area = img.width * img.height
                    if area > best_area and area > 10000:
                        best_area = area
                        best = (img_bytes, shape.image.ext)
                except Exception:
                    pass

    if best:
        img_bytes, img_ext = best
        out_path = tmp_dir / f"{prefix}_produit.{img_ext}"
        with open(out_path, "wb") as f:
            f.write(img_bytes)
        return out_path

    return None


# ════════════════════════════════════════════════════════════════
#  4. EXTRACTION DIAGRAMME POLAIRE
# ════════════════════════════════════════════════════════════════

def extract_polaire(input_path: Path, tmp_dir: Path, has_polaire: bool) -> Path | None:
    """Extrait le diagramme polaire si disponible."""
    if not has_polaire:
        return None

    ext = input_path.suffix.lower()

    if ext == ".pdf":
        return _find_polaire_in_pdf(input_path, tmp_dir)
    elif ext in (".pptx", ".ppt"):
        return _find_polaire_in_pptx(input_path, tmp_dir)
    elif ext in (".jpg", ".jpeg", ".png"):
        # Détecte et recadre la zone circulaire (diagramme polaire)
        return _detect_polar_region(input_path, tmp_dir)

    return None


def _find_polaire_in_pdf(pdf_path: Path, tmp_dir: Path) -> Path | None:
    """
    Cherche le diagramme polaire dans un PDF.
    Heuristique : image carrée ou quasi-carrée (ratio 0.7-1.3) et
    de taille moyenne (pas la plus grande = photo produit).
    """
    doc    = fitz.open(str(pdf_path))
    images = []

    for page in doc:
        for img_info in page.get_images(full=True):
            xref     = img_info[0]
            base_img = doc.extract_image(xref)
            w, h     = base_img["width"], base_img["height"]
            area     = w * h
            if area < 5000:
                continue
            ratio = w / h if h > 0 else 0
            # Diagramme polaire : quasi-carré
            if 0.7 <= ratio <= 1.4:
                images.append((area, xref, base_img["ext"]))

    doc_reload = fitz.open(str(pdf_path))
    # Trie par taille, prend le 2ème plus grand (le plus grand = photo produit)
    images.sort(reverse=True)
    candidates = images[1:] if len(images) > 1 else images[:1]

    for _, xref, img_ext in candidates[:1]:
        base_img = doc_reload.extract_image(xref)
        out_path = tmp_dir / f"polaire.{img_ext}"
        with open(out_path, "wb") as f:
            f.write(base_img["image"])
        doc_reload.close()
        return out_path

    doc_reload.close()
    return None


def _find_polaire_in_pptx(pptx_path: Path, tmp_dir: Path) -> Path | None:
    """Cherche le diagramme polaire dans un PPTX (même heuristique)."""
    from pptx import Presentation as PptxPrs
    prs     = PptxPrs(str(pptx_path))
    images  = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                try:
                    img = Image.open(__import__("io").BytesIO(shape.image.blob))
                    w, h  = img.width, img.height
                    ratio = w / h if h > 0 else 0
                    area  = w * h
                    if area > 5000 and 0.7 <= ratio <= 1.4:
                        images.append((area, shape.image.blob, shape.image.ext))
                except Exception:
                    pass

    images.sort(reverse=True)
    candidates = images[1:] if len(images) > 1 else images[:1]
    for _, blob, ext in candidates[:1]:
        out_path = tmp_dir / f"polaire.{ext}"
        with open(out_path, "wb") as f:
            f.write(blob)
        return out_path

    return None


def _detect_polar_region(img_path: Path, tmp_dir: Path) -> Path | None:
    """Détecte la zone circulaire d'un diagramme polaire dans une image raster."""
    try:
        img = Image.open(img_path).convert("RGB")
        arr = np.array(img)
        # Cherche les pixels sombres (courbes photométriques)
        dark = np.all(arr < 100, axis=2)
        rows = np.any(dark, axis=1)
        cols = np.any(dark, axis=0)
        if not rows.any():
            return None
        y1, y2 = np.where(rows)[0][[0, -1]]
        x1, x2 = np.where(cols)[0][[0, -1]]
        pad = 20
        region = (max(0, x1-pad), max(0, y1-pad),
                  min(img.width, x2+pad), min(img.height, y2+pad))
        crop = img.crop(region)
        out  = tmp_dir / "polaire.png"
        crop.save(str(out))
        return out
    except Exception:
        return None


# ════════════════════════════════════════════════════════════════
#  5. GÉNÉRATION PPTX
# ════════════════════════════════════════════════════════════════

def generate_pptx(cfg: dict, output_path: Path):
    """Écrit le JSON de config et appelle generate_pptx.js."""
    cfg_path = output_path.parent / "config_fiche.json"
    with open(cfg_path, "w", encoding="utf-8") as f:
        json.dump(cfg, f, ensure_ascii=False, indent=2)

    result = subprocess.run(
        ["node", str(GENERATE_JS), str(cfg_path), str(output_path)],
        capture_output=True, text=True
    )
    if result.returncode != 0:
        raise RuntimeError(f"Erreur generate_pptx.js :\n{result.stderr}")
    print(result.stdout.strip())


# ════════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════════

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    input_path  = Path(sys.argv[1]).resolve()
    output_name = sys.argv[2] if len(sys.argv) > 2 else None

    if not input_path.exists():
        print(f"[ERREUR] Fichier introuvable : {input_path}")
        sys.exit(1)

    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("[ERREUR] Variable ANTHROPIC_API_KEY non définie.")
        print("  export ANTHROPIC_API_KEY=sk-ant-...")
        sys.exit(1)

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        print(f"[1/5] Conversion en images…  ({input_path.name})")
        images = file_to_images(input_path, tmp_dir)
        print(f"      → {len(images)} page(s) détectée(s)")

        print("[2/5] Extraction des données produit via Claude API…")
        data = extract_product_data(images)
        nom  = data.get("nom_produit", "PRODUIT")
        print(f"      → {nom}")

        print("[3/5] Extraction photo principale…")
        photo_path = extract_main_photo(input_path, tmp_dir)
        if photo_path:
            print(f"      → {photo_path.name}")
        else:
            print("      → aucune photo trouvée")

        print("[4/5] Extraction diagramme polaire…")
        polaire_path = extract_polaire(
            input_path, tmp_dir, data.get("a_diagramme_polaire", False)
        )
        if polaire_path:
            print(f"      → {polaire_path.name}")
        else:
            print("      → non disponible")

        # Construction config finale
        ref_clean = (data.get("reference") or nom[:12]).replace(" ", "-").upper()
        annee     = str(data.get("annee") or __import__("datetime").date.today().year)

        if output_name is None:
            output_name = f"FT-{ref_clean}-{annee}.pptx"

        output_path = Path(output_name).resolve()

        # Specs : filtre les valeurs null
        specs_clean = [
            [lbl, val] for lbl, val in (data.get("specifications") or [])
            if val and val != "null"
        ]

        cfg = {
            "nom_produit"   : data.get("nom_produit", "NOM DU PRODUIT"),
            "reference"     : ref_clean,
            "serie"         : data.get("serie", "Éclairage industriel"),
            "gamme"         : data.get("gamme"),
            "accroche"      : data.get("accroche", ""),
            "description"   : data.get("description", ""),
            "specifications": specs_clean,
            "certifications": data.get("certifications", {}),
            "dimensions"    : data.get("dimensions", []),
            "photo_path"    : str(photo_path) if photo_path else None,
            "polaire_path"  : str(polaire_path) if polaire_path else None,
            "schema_path"   : None,   # extrait séparément si besoin
            "annee"         : annee,
        }

        # Copie les images extraites dans le répertoire de sortie
        # (pour que Node.js puisse les lire après fermeture du tempdir)
        out_dir = output_path.parent
        if photo_path and photo_path.exists():
            dst = out_dir / photo_path.name
            shutil.copy(photo_path, dst)
            cfg["photo_path"] = str(dst)
        if polaire_path and polaire_path.exists():
            dst = out_dir / polaire_path.name
            shutil.copy(polaire_path, dst)
            cfg["polaire_path"] = str(dst)

        print(f"[5/5] Génération PPTX…  → {output_path.name}")
        generate_pptx(cfg, output_path)

    print(f"\n✅  Fiche technique générée : {output_path}")


if __name__ == "__main__":
    main()
